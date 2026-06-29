from __future__ import annotations

import re
from pathlib import Path
from typing import List


def parse_timeline(md_text: str):
    events = []
    # Split by lines and find event headers starting with '- **'
    lines = md_text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if line.startswith('- **') and '—' in line:
            # Extract title between '**' and '**' or full line
            m = re.match(r"- \*\*(\d{4}) — (.+?)\*\*", line)
            if m:
                year = m.group(1)
                title = m.group(2).strip()
            else:
                # fallback parse
                parts = line[3:].split('—', 1)
                year = parts[0].strip().strip('* ')
                title = parts[1].strip().strip('* ')
            # collect following indented lines (Influence / Notes)
            influence = ''
            notes = ''
            i += 1
            while i < len(lines) and lines[i].strip().startswith('-'):
                sub = lines[i].strip()
                if sub.lower().startswith('- influence:'):
                    influence = sub.split(':', 1)[1].strip()
                elif sub.lower().startswith('- notes:'):
                    notes = sub.split(':', 1)[1].strip()
                i += 1
            events.append({'year': year, 'title': title, 'influence': influence, 'notes': notes})
        else:
            i += 1
    return events


def make_pptx(events, out_path: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    from pptx.dml.color import RGBColor

    # Timing defaults (minutes)
    TITLE_SLIDE_MINUTES = 1
    DEFAULT_SLIDE_MINUTES = 2

    def fmt_time(minutes_total: float):
        secs = int(minutes_total * 60)
        m = secs // 60
        s = secs % 60
        return f"{m:02d}:{s:02d}"

    prs = Presentation()
    # title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Turring Influence Timeline"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Key events and their influence on AI"
    bullet_slide_layout = prs.slide_layouts[1]
    total_slides = 1 + len(events)

    # style title slide background and footer
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x1F, 0x4A, 0x78)
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    except Exception:
        # best-effort styling; proceed if slide master doesn't allow
        pass

    tb = slide.shapes.add_textbox(Inches(0.3), Inches(6.6), Inches(6), Inches(0.4))
    tf = tb.text_frame
    tf.text = 'Turring — For internal review'
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(10)

    # add slide number and est time on title
    sn_box = slide.shapes.add_textbox(Inches(7.2), Inches(6.6), Inches(2), Inches(0.4))
    sn_box.text = f"Slide 1 of {total_slides} — Est. {TITLE_SLIDE_MINUTES} min"

    # cumulative timing in minutes
    cumulative = TITLE_SLIDE_MINUTES
    for idx, ev in enumerate(events, start=2):
        slide = prs.slides.add_slide(bullet_slide_layout)
        # per-slide background tint (subtle)
        try:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFB)
        except Exception:
            pass

        slide.shapes.title.text = f"{ev['year']}: {ev['title']}"
        body = slide.shapes.placeholders[1].text_frame
        # clear default paragraph
        if body.paragraphs:
            body.paragraphs[0].text = ''

        # Add concise bullets (summary first, then significance)
        if ev.get('influence'):
            p = body.add_paragraph()
            p.text = ev['influence']
            p.level = 0
        if ev.get('notes'):
            p = body.add_paragraph()
            p.text = ev['notes']
            p.level = 1

        # Style paragraphs
        for paragraph in body.paragraphs:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(14)

        # Footer (left) and slide number + est time (right)
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(6.6), Inches(5), Inches(0.4))
        footer_tf = footer.text_frame
        footer_tf.text = 'Turring — Prepared for coach review'
        for p in footer_tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(9)

        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(6.6), Inches(2), Inches(0.4))
        right_box.text = f"Slide {idx} of {total_slides} — Est. {DEFAULT_SLIDE_MINUTES} min"
        for p in right_box.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(9)
                r.italic = True

        # Speaker notes with timestamps to align student outputs and lesson throughput
        start_min = cumulative
        end_min = cumulative + DEFAULT_SLIDE_MINUTES
        student_demo_min = start_min + 0.5 * DEFAULT_SLIDE_MINUTES
        # format times
        ts_start = fmt_time(start_min)
        ts_demo = fmt_time(student_demo_min)
        ts_qa = fmt_time(end_min)

        notes_lines = []
        notes_lines.append(f"Overview: {ev['year']} — {ev['title']}")
        if ev.get('influence'):
            notes_lines.append(f"Key influence: {ev['influence']}")
        if ev.get('notes'):
            notes_lines.append(f"Details: {ev['notes']}")
        notes_lines.append("")
        notes_lines.append("Timestamps (mm:ss):")
        notes_lines.append(f"- Start: {ts_start}")
        notes_lines.append(f"- Student output / demo: {ts_demo}")
        notes_lines.append(f"- Q&A / discussion: {ts_qa}")
        notes_lines.append("")
        notes_lines.append("Suggested talking points:")
        notes_lines.append("- One-line summary of impact")
        notes_lines.append("- How this ties to the model's outputs (show example)")
        notes_lines.append("- 1 question for coach/team")

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = '\n'.join(notes_lines)

        cumulative = end_min

    prs.save(str(out_path))


def make_docx(events, out_path: Path):
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading('Turring Influence Timeline', 0)
    doc.add_paragraph('Key events and their influence on AI')

    for ev in events:
        doc.add_heading(f"{ev['year']}: {ev['title']}", level=2)
        if ev.get('influence'):
            p = doc.add_paragraph()
            run = p.add_run('Influence: ' + ev['influence'])
            run.font.size = Pt(11)
        if ev.get('notes'):
            p = doc.add_paragraph()
            run = p.add_run('Notes: ' + ev['notes'])
            run.font.size = Pt(11)

        # Speaker notes
        sp = doc.add_paragraph()
        sp_run = sp.add_run('Speaker notes:')
        sp_run.bold = True
        sp_run.font.size = Pt(10)
        sn = doc.add_paragraph()
        sn_lines = []
        sn_lines.append(f"Overview: {ev['year']} — {ev['title']}")
        if ev.get('influence'):
            sn_lines.append(f"Key influence: {ev['influence']}")
        if ev.get('notes'):
            sn_lines.append(f"Details: {ev['notes']}")
        sn_lines.append('')
        sn_lines.append('Suggested talking points:')
        sn_lines.append('- One-sentence summary of impact')
        sn_lines.append('- Connection to team aims or model outputs')
        sn_run = sn.add_run('\n'.join(sn_lines))
        sn_run.font.size = Pt(10)

    doc.save(str(out_path))


def main():
    root = Path(__file__).resolve().parent
    md_file = root / 'timeline.md'
    if not md_file.exists():
        print('timeline.md not found')
        return
    md_text = md_file.read_text(encoding='utf-8')
    events = parse_timeline(md_text)
    if not events:
        print('No events parsed from timeline.md')
        return

    pptx_path = root / 'Turring_Timeline.pptx'
    docx_path = root / 'Turring_Timeline.docx'

    make_pptx(events, pptx_path)
    make_docx(events, docx_path)

    print(f'Wrote {pptx_path}')
    print(f'Wrote {docx_path}')


if __name__ == '__main__':
    main()
